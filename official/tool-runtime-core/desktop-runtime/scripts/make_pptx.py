#!/usr/bin/env python
"""
make_pptx.py — 直接用 python-pptx 生成专业演示文稿
输入：JSON 文件路径（设计规格）
输出：PPTX 文件路径（打印到 stdout）
"""
import sys, json, os
from pathlib import Path

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu, Cm
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    from lxml import etree
except ImportError as e:
    print(json.dumps({"ok": False, "error": str(e)}, ensure_ascii=False))
    sys.exit(1)

# ── 单位换算 ──────────────────────────────────────────────────────
W = Inches(13.333)   # 16:9 宽
H = Inches(7.5)      # 16:9 高

def rgb(hex_str):
    h = hex_str.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def add_rect(slide, x, y, w, h_px, fill_hex, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h_px)
    shape.line.fill.background()
    shape.line.width = 0
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = rgb(fill_hex)
    if alpha is not None:
        fill.fore_color.theme_color  # sometimes needed
        # set alpha via xml
        sp = shape.element
        solidFill = sp.find('.//' + qn('a:solidFill'))
        if solidFill is not None:
            srgb = solidFill.find(qn('a:srgbClr'))
            if srgb is not None:
                alpha_el = etree.SubElement(srgb, qn('a:alpha'))
                alpha_el.set('val', str(int(alpha * 100000)))
    return shape

def add_textbox(slide, x, y, w, h_px, text, font_size, font_color, bold=False,
                align=PP_ALIGN.LEFT, font_name='Microsoft YaHei', word_wrap=True):
    if not text:
        return None
    tb = slide.shapes.add_textbox(x, y, w, h_px)
    tf = tb.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = rgb(font_color)
    run.font.bold = bold
    run.font.name = font_name
    return tb

def add_text_lines(slide, x, y, w, h_px, lines, font_size, font_color,
                   bold=False, align=PP_ALIGN.LEFT, line_space=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h_px)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = str(line)
        run.font.size = Pt(font_size)
        run.font.color.rgb = rgb(font_color)
        run.font.bold = bold
        run.font.name = 'Microsoft YaHei'
        # line spacing
        from pptx.oxml.ns import nsmap
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPct = etree.SubElement(lnSpc, qn('a:spcPct'))
        spcPct.set('val', f'{int(line_space*100000)}')
    return tb

# ── 幻灯片生成函数 ────────────────────────────────────────────────

def make_cover(prs, slide_spec, colors):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    primary = colors['primary']
    accent  = colors['accent']
    light   = colors['text_light']
    muted   = colors.get('text_muted', '#AACCEE')

    add_rect(slide, 0, 0, W, H, primary)
    add_rect(slide, 0, 0, W, Inches(0.08), accent)
    add_rect(slide, 0, H - Inches(1.4), W, Inches(1.4), '#0d2240')
    add_rect(slide, 0, 0, Inches(0.12), H, accent)
    s = slide.shapes.add_shape(9, W - Inches(3.2), Inches(0.3), Inches(3.5), Inches(3.5))
    s.line.fill.background(); s.line.width = 0
    s.fill.solid(); s.fill.fore_color.rgb = rgb('#0d2240')

    title = slide_spec.get('title', '')
    subtitle = slide_spec.get('subtitle', '')
    add_textbox(slide, Inches(0.4), Inches(1.8), Inches(9), Inches(1.6),
                title, 36, light, bold=True, align=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide, Inches(0.4), Inches(3.5), Inches(9), Inches(0.8),
                    subtitle, 20, accent, bold=False)
    import datetime
    date_str = datetime.date.today().strftime('%Y年%m月')
    add_textbox(slide, Inches(0.4), H - Inches(1.1), Inches(6), Inches(0.6),
                date_str, 13, muted, bold=False)


def make_toc(prs, slide_spec, colors):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    primary = colors['primary']
    accent  = colors['accent']
    light   = colors['text_light']
    bg_c    = colors.get('bg_content', '#F0F4F8')

    add_rect(slide, 0, 0, Inches(3.8), H, primary)
    add_rect(slide, 0, 0, Inches(0.12), H, accent)
    add_textbox(slide, Inches(0.25), Inches(2.5), Inches(3.3), Inches(1.2),
                slide_spec.get('title', '目录'), 28, light, bold=True, align=PP_ALIGN.LEFT)
    add_rect(slide, Inches(3.8), 0, W - Inches(3.8), H, bg_c)

    items = slide_spec.get('items', [])
    for i, item in enumerate(items[:6]):
        y = Inches(0.7 + i * 0.95)
        x_num = Inches(4.1)
        x_txt = Inches(5.0)
        num_bg = slide.shapes.add_shape(9, x_num, y + Pt(2), Inches(0.55), Inches(0.55))
        num_bg.line.fill.background(); num_bg.line.width = 0
        num_bg.fill.solid(); num_bg.fill.fore_color.rgb = rgb(accent)
        add_textbox(slide, x_num, y, Inches(0.55), Inches(0.55),
                    f'{i+1:02d}', 18, light, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, x_txt, y + Pt(4), Inches(8.0), Inches(0.55),
                    str(item), 16, primary, bold=False)
        if i < len(items) - 1:
            sep = slide.shapes.add_shape(1, x_txt, y + Inches(0.72), Inches(7.8), Pt(1))
            sep.line.fill.background(); sep.line.width = 0
            sep.fill.solid(); sep.fill.fore_color.rgb = rgb('#CCDDEE')

    _add_page_num(slide, len(prs.slides), light, primary)


def make_chapter(prs, slide_spec, colors, chapter_num=1):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    primary = colors['primary']
    accent  = colors['accent']
    light   = colors['text_light']
    bg_c    = colors.get('bg_content', '#F0F4F8')

    add_rect(slide, 0, 0, Inches(3.8), H, primary)
    add_rect(slide, 0, 0, Inches(0.12), H, accent)
    add_rect(slide, Inches(3.8), 0, W - Inches(3.8), H, bg_c)
    add_textbox(slide, Inches(0.3), Inches(1.8), Inches(3.2), Inches(1.2),
                f'{chapter_num:02d}', 52, accent, bold=True, align=PP_ALIGN.LEFT)
    add_textbox(slide, Inches(0.3), Inches(3.0), Inches(3.2), Inches(1.5),
                slide_spec.get('title', ''), 22, light, bold=True, align=PP_ALIGN.LEFT)
    _add_page_num(slide, len(prs.slides), light, primary)


def make_content(prs, slide_spec, colors):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    primary = colors['primary']
    accent  = colors['accent']
    light   = colors['text_light']
    dark    = colors.get('text_dark', '#222222')
    bg_c    = colors.get('bg_content', '#F0F4F8')

    # 左侧深蓝标题栏
    add_rect(slide, 0, 0, Inches(3.2), H, primary)
    add_rect(slide, 0, 0, Inches(0.12), H, accent)
    # 右侧浅色内容区
    add_rect(slide, Inches(3.2), 0, W - Inches(3.2), H, bg_c)

    title = slide_spec.get('title', '')
    add_textbox(slide, Inches(0.25), Inches(0.4), Inches(2.8), Inches(5.5),
                title, 20, light, bold=True, align=PP_ALIGN.LEFT, word_wrap=True)

    points = slide_spec.get('points', [])
    if not points and 'content' in slide_spec:
        points = [slide_spec['content']] if isinstance(slide_spec['content'], str) else slide_spec['content']

    n = max(len(points), 1)
    available_h = Inches(6.2)
    card_h = min(available_h / n, Inches(1.35))
    font_size = 15 if n <= 2 else (13 if n == 3 else 12)
    gap = Inches(0.1)

    for i, pt in enumerate(points):
        y = Inches(0.4) + i * (card_h + gap)
        x = Inches(3.4)
        w_card = W - Inches(3.6)
        bar = slide.shapes.add_shape(1, x, y + Inches(0.08), Pt(4), card_h - Inches(0.16))
        bar.line.fill.background(); bar.line.width = 0
        bar.fill.solid(); bar.fill.fore_color.rgb = rgb(accent)
        add_textbox(slide, x + Inches(0.15), y + Pt(6), w_card - Inches(0.2),
                    card_h - Pt(10), str(pt), font_size, dark, word_wrap=True)
        if i < len(points) - 1:
            sep = slide.shapes.add_shape(1, x, y + card_h + Inches(0.04), w_card, Pt(1))
            sep.line.fill.background(); sep.line.width = 0
            sep.fill.solid(); sep.fill.fore_color.rgb = rgb('#CCDDEE')

    _add_page_num(slide, len(prs.slides), light, primary)


def make_ending(prs, slide_spec, colors):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    primary = colors['primary']
    accent  = colors['accent']
    light   = colors['text_light']

    add_rect(slide, 0, 0, W, H, primary)
    add_rect(slide, 0, 0, Inches(0.12), H, accent)
    add_rect(slide, 0, 0, W, Inches(0.08), accent)
    add_rect(slide, 0, H - Inches(0.08), W, Inches(0.08), accent)
    s = slide.shapes.add_shape(9, W - Inches(3.0), Inches(0.5), Inches(3.2), Inches(3.2))
    s.line.fill.background(); s.line.width = 0
    s.fill.solid(); s.fill.fore_color.rgb = rgb('#0d2240')

    add_textbox(slide, Inches(0.5), Inches(2.0), Inches(11), Inches(1.8),
                slide_spec.get('title', '感谢聆听'), 52, light, bold=True, align=PP_ALIGN.CENTER)
    msg = slide_spec.get('message', '')
    if msg:
        add_textbox(slide, Inches(0.5), Inches(3.9), Inches(11), Inches(0.9),
                    msg, 20, accent, align=PP_ALIGN.CENTER)
    _add_page_num(slide, len(prs.slides), light, primary)


def _add_page_num(slide, num, light_color, primary_color):
    add_textbox(slide, W - Inches(0.6), H - Inches(0.4), Inches(0.5), Inches(0.35),
                str(num), 11, light_color, bold=True, align=PP_ALIGN.CENTER)


# ── 主流程 ────────────────────────────────────────────────────────
def main():
    if len(sys.argv) < 3:
        print(json.dumps({'ok': False, 'error': 'usage: make_pptx.py <spec.json> <output.pptx>'}, ensure_ascii=False))
        sys.exit(1)

    spec_path   = sys.argv[1]
    output_path = sys.argv[2]

    with open(spec_path, encoding='utf-8-sig') as f:
        spec = json.load(f)

    colors = spec.get('color_scheme', {
        'primary':    '#1A3A5C',
        'secondary':  '#2d5f8a',
        'accent':     '#E8A020',
        'background': '#1A3A5C',
        'bg_content': '#F0F4F8',
        'text_dark':  '#222222',
        'text_light': '#FFFFFF',
        'text_muted': '#AACCEE',
    })
    colors.setdefault('primary',    '#1A3A5C')
    colors.setdefault('secondary',  '#2d5f8a')
    colors.setdefault('accent',     '#E8A020')
    colors.setdefault('background', '#1A3A5C')
    colors.setdefault('bg_content', '#F0F4F8')
    colors.setdefault('text_dark',  '#222222')
    colors.setdefault('text_light', '#FFFFFF')
    colors.setdefault('text_muted', '#AACCEE')

    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    chapter_count = 0
    for slide_spec in spec.get('slides', []):
        stype = slide_spec.get('type', 'content')
        if stype == 'cover':
            make_cover(prs, slide_spec, colors)
        elif stype == 'toc':
            make_toc(prs, slide_spec, colors)
        elif stype == 'chapter':
            chapter_count += 1
            make_chapter(prs, slide_spec, colors, chapter_count)
        elif stype == 'ending':
            make_ending(prs, slide_spec, colors)
        else:
            make_content(prs, slide_spec, colors)

    prs.save(output_path)
    print(json.dumps({'ok': True, 'path': output_path, 'slides': len(prs.slides)}, ensure_ascii=False))

if __name__ == '__main__':
    main()
